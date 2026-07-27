"""
make_pptx.py — Converts presentation content into presentation.pptx
using python-pptx with a sleek dark modern aesthetic matching presentation.html.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG = RGBColor(11, 15, 25)        # #0b0f19
    COLOR_CARD = RGBColor(18, 26, 44)      # #121a2c
    COLOR_CARD_BORDER = RGBColor(30, 41, 59) # #1e293b
    COLOR_TEAL = RGBColor(20, 184, 166)    # #14b8a6
    COLOR_CYAN = RGBColor(6, 182, 212)     # #06b6d4
    COLOR_PURPLE = RGBColor(139, 92, 246)  # #8b5cf6
    COLOR_EMERALD = RGBColor(16, 185, 129) # #10b981
    COLOR_TEXT_MAIN = RGBColor(243, 244, 246) # #f3f4f6
    COLOR_TEXT_MUTED = RGBColor(156, 163, 175) # #9ca3af
    COLOR_TEXT_DIM = RGBColor(107, 114, 128) # #6b7280

    blank_slide_layout = prs.slide_layouts[6] # blank layout

    def add_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, tag_text, title_text, subtitle_text):
        add_bg(slide)

        # Slide Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_TEAL
        p_tag.font.name = 'Arial'

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN
        p_title.font.name = 'Arial'

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.6))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED
        p_sub.font.name = 'Arial'

    def add_card(slide, left, top, width, height, title, content_list, title_color=COLOR_TEXT_MAIN):
        # Card container shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_CARD_BORDER
        shape.line.width = Pt(1)

        # Card Title
        tb_title = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_t = tf_title.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = title_color
        p_t.font.name = 'Arial'

        # Card Content
        tb_content = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.65), width - Inches(0.4), height - Inches(0.8))
        tf_c = tb_content.text_frame
        tf_c.word_wrap = True
        
        for idx, text in enumerate(content_list):
            p = tf_c.paragraphs[0] if idx == 0 else tf_c.add_paragraph()
            p.text = text
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_TEXT_MUTED
            p.font.name = 'Arial'
            p.space_after = Pt(8)

    # ── SLIDE 1: Title Slide ───────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide1)

    # Icon/Emoji
    tb_icon = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.8))
    p_icon = tb_icon.text_frame.paragraphs[0]
    p_icon.text = "🧬"
    p_icon.alignment = PP_ALIGN.CENTER
    p_icon.font.size = Pt(48)

    # Main Title
    tb_t1 = slide1.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.733), Inches(1.2))
    p_t1 = tb_t1.text_frame.paragraphs[0]
    p_t1.text = "GEO Meta-Analysis Toolkit"
    p_t1.alignment = PP_ALIGN.CENTER
    p_t1.font.size = Pt(44)
    p_t1.font.bold = True
    p_t1.font.color.rgb = COLOR_TEXT_MAIN

    # Subtitle
    tb_s1 = slide1.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.333), Inches(1.0))
    p_s1 = tb_s1.text_frame.paragraphs[0]
    p_s1.text = "Autonomous Multi-Tool Agent for NCBI GEO Exploration, Metadata Aggregation, & AI Narrative Landscape Synthesis"
    p_s1.alignment = PP_ALIGN.CENTER
    p_s1.font.size = Pt(18)
    p_s1.font.color.rgb = COLOR_TEXT_MUTED

    # Highlights Badge Row
    tb_b1 = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.0))
    tf_b1 = tb_b1.text_frame
    p_b1 = tf_b1.paragraphs[0]
    p_b1.text = "⚡ 3-Phase Smart Hybrid Engine   |   🧠 Local Ollama medgemma:4b   |   🔬 NCBI Entrez E-Utils   |   🧬 FAISS Vector Index"
    p_b1.alignment = PP_ALIGN.CENTER
    p_b1.font.size = Pt(14)
    p_b1.font.bold = True
    p_b1.font.color.rgb = COLOR_TEAL

    slide1.notes_slide.notes_text_frame.text = "Welcome everyone! Today we present the GEO Meta-Analysis Toolkit — an autonomous, self-healing biomedical research agent engineered to query, aggregate, and synthesize gene expression datasets from NCBI GEO effortlessly."


    # ── SLIDE 2: Problem Statement ─────────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "The Challenge", "Biomedical Discovery Bottlenecks", "Why traditional manual search and generic linear scripts fall short in high-throughput genomics.")

    add_card(slide2, Inches(0.8), Inches(2.2), Inches(3.6), Inches(4.5), 
             "📦 Data Scale & Complexity", 
             ["• NCBI GEO holds >160,000 datasets across thousands of organisms.",
              "• Complex search syntax (Entrez Boolean & Field tags).",
              "• Non-standardized sample descriptions and metadata schemas."])

    add_card(slide2, Inches(4.8), Inches(2.2), Inches(3.6), Inches(4.5), 
             "⏳ Manual Synthesis Time", 
             ["• Extracting organism distributions and timelines takes hours.",
              "• Manual cross-referencing of PubMed literature.",
              "• Prone to human error when handling hundreds of GSE records."])

    add_card(slide2, Inches(8.8), Inches(2.2), Inches(3.6), Inches(4.5), 
             "🤖 LLM & Hardware Limits", 
             ["• Pure agentic ReAct loops stall on LLM reasoning timeouts.",
              "• Rate limits and API cost overruns.",
              "• Laptops constrained by GPU VRAM (e.g. 8 GB VRAM GPU)."])

    slide2.notes_slide.notes_text_frame.text = "Researchers face a daunting challenge: NCBI GEO hosts over 160,000 datasets. Manual discovery takes hours, metadata formats are inconsistent, literature cross-referencing is fragmented, and classical AI tools suffer from rate limits or timeouts."


    # ── SLIDE 3: The Solution ──────────────────────────────────────────────────
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "The Solution", "Unified AI-Powered Meta-Analysis", "Bridging deterministic high-speed computation with autonomous agent reasoning.")

    stat_data = [
        ("< 2s", "Instant Aggregation", "Phase 1 direct parsing provides real-time chart rendering before LLM starts."),
        ("100%", "Execution Resilience", "Self-healing 3-phase execution ensures no search query ever fails."),
        ("3.5 GB", "VRAM Efficiency", "Tailored for local GPU acceleration using bio LLM medgemma:4b."),
        ("10,000", "Max Scalability", "Scalable data fetching engine tuned with NCBI API key rate-limiting.")
    ]

    for idx, (num, label, desc) in enumerate(stat_data):
        left = Inches(0.8 + idx * 3.0)
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.4), Inches(2.7), Inches(4.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_CARD_BORDER
        
        # Num
        tb = slide3.shapes.add_textbox(left + Inches(0.1), Inches(2.6), Inches(2.5), Inches(0.9))
        p = tb.text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEAL

        # Label
        tb_l = slide3.shapes.add_textbox(left + Inches(0.1), Inches(3.6), Inches(2.5), Inches(0.5))
        p_l = tb_l.text_frame.paragraphs[0]
        p_l.text = label.upper()
        p_l.font.size = Pt(11)
        p_l.font.bold = True
        p_l.font.color.rgb = COLOR_TEXT_MAIN

        # Desc
        tb_d = slide3.shapes.add_textbox(left + Inches(0.1), Inches(4.2), Inches(2.5), Inches(2.2))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_TEXT_MUTED

    slide3.notes_slide.notes_text_frame.text = "Our solution is the GEO Meta-Analysis Toolkit: a production-grade Streamlit application driven by a 3-Phase Unified Smart Hybrid Engine that guarantees 100% execution success while leveraging local GPU AI."


    # ── SLIDE 4: System Architecture ───────────────────────────────────────────
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "System Architecture", "Unified Smart Hybrid Architecture", "Modular design decoupling user interaction, agent reasoning, API tools, and semantic RAG vector search.")

    # UI Box
    add_card(slide4, Inches(2.6), Inches(2.2), Inches(8.1), Inches(1.1), "Streamlit Dark UI Dashboard", ["Summary Metrics • Visual Charts • Datasets Table • Reasoning Trace"], COLOR_CYAN)
    
    # Engine Box
    add_card(slide4, Inches(1.6), Inches(3.6), Inches(10.1), Inches(1.3), "Unified Smart Hybrid Engine", ["Phase 1: Direct GEO E-Utils → Phase 2: LangChain ReAct PubMed → Phase 3: Auto Fallback"], COLOR_TEAL)

    # Bottom 3 Boxes
    add_card(slide4, Inches(0.8), Inches(5.2), Inches(3.6), Inches(1.6), "Local GPU LLM", ["Ollama (medgemma:4b)", "Strict Local GPU/CPU Execution"])
    add_card(slide4, Inches(4.8), Inches(5.2), Inches(3.6), Inches(1.6), "NCBI E-Utilities", ["GEO Datasets Search", "PubMed Abstract Retrieval"])
    add_card(slide4, Inches(8.8), Inches(5.2), Inches(3.6), Inches(1.6), "FAISS Vector Store", ["Cell Ontology (CL, UBERON, EFO)", "Two-Pass LLM Normalizer"])

    slide4.notes_slide.notes_text_frame.text = "Here is the architecture of the system: Streamlit UI interfaces with the Unified Smart Hybrid Engine, which manages Phase 1 Direct GEO query, Phase 2 LangChain ReAct PubMed investigation, and Phase 3 linear fallback."


    # ── SLIDE 5: 3-Phase Execution Workflow ────────────────────────────────────
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "Execution Strategy", "Self-Healing 3-Phase Workflow", "Eliminating LLM single-point failures through layered execution fallback.")

    add_card(slide5, Inches(0.8), Inches(2.2), Inches(3.6), Inches(4.5), 
             "Phase 1: Direct Engine", 
             ["⚡ Execution Time: ~2 seconds",
              "• Queries NCBI GEO directly via E-utilities API.",
              "• Parses XML metadata instantly.",
              "• Extracts organism counts, timelines, sample stats, and title terms."], 
             COLOR_CYAN)

    add_card(slide5, Inches(4.8), Inches(2.2), Inches(3.6), Inches(4.5), 
             "Phase 2: ReAct AI Synthesis", 
             ["🧠 Execution Time: ~10-15 seconds",
              "• Spawns specialized LangChain ReAct Agent.",
              "• Autonomously queries PubMed for literature context.",
              "• Synthesizes qualitative research landscape summary."], 
             COLOR_PURPLE)

    add_card(slide5, Inches(8.8), Inches(2.2), Inches(3.6), Inches(4.5), 
             "Phase 3: Automatic Fallback", 
             ["🛡️ 100% Uptime Shield",
              "• Catches LLM timeouts, parsing errors, or network issues.",
              "• Seamlessly auto-degrades to linear summarizer.",
              "• Guarantees zero empty outputs on any search query."], 
             COLOR_EMERALD)

    slide5.notes_slide.notes_text_frame.text = "The core innovation is our 3-phase execution model. Phase 1 provides instant deterministic metadata statistics in ~2s. Phase 2 spawns a LangChain ReAct agent to investigate PubMed literature. If any step times out, Phase 3 seamlessly auto-degrades to the linear summarizer."


    # ── SLIDE 6: Ontology Disambiguation & RAG ─────────────────────────────────
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide6, "Semantic Intelligence", "Ontology Disambiguation & RAG", "Standardizing fragmented biomedical jargon using FAISS vector search & ontology mappings.")

    add_card(slide6, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5), 
             "📚 Supported Biomedical Ontologies", 
             ["• Cell Ontology (CL): Standardized cell types (e.g., CD8+ T cell, Podocyte, Microglia).",
              "• UBERON: Anatomical structures and tissue locations across species.",
              "• Experimental Factor Ontology (EFO): Experimental assay types, variables, and platform technology."], 
             COLOR_TEAL)

    add_card(slide6, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.5), 
             "💡 Two-Pass Cost-Gated Disambiguation", 
             ["1. Fast Vector Search: Cosine similarity lookup in pre-computed FAISS index.",
              "2. Cost-Gated LLM Normalizer: Triggers ONLY when similarity is below 0.70 threshold.",
              "• Resolves ambiguous jargon (e.g., 'MACS' cell separation method vs. 'macs' slang for macrophages).",
              "• Zero extra LLM round-trips for clean terms."], 
             COLOR_CYAN)

    slide6.notes_slide.notes_text_frame.text = "We also integrate a semantic ontology vector store using FAISS and sentence-transformers. It handles ambiguous biomedical terms like 'MACS' using a cost-gated two-pass LLM normalizer."


    # ── SLIDE 7: Dashboard Views ──────────────────────────────────────────────
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide7, "Interactive Analytics", "Multi-Tab Meta-Analysis Dashboard", "Live aggregated analytics across NCBI GEO dataset search results.")

    add_card(slide7, Inches(0.8), Inches(2.2), Inches(5.6), Inches(2.1), "📝 Summary Tab", ["• Key summary metrics (total datasets, total sample volume, top organism).", "• Flowing LLM narrative overview of the research landscape."])
    add_card(slide7, Inches(6.8), Inches(2.2), Inches(5.7), Inches(2.1), "📈 Visual Charts Tab", ["• Organism distribution donut chart.", "• Yearly publication timeline & top title keyword term frequency."])
    add_card(slide7, Inches(0.8), Inches(4.6), Inches(5.6), Inches(2.1), "📋 Datasets Table", ["• Full sortable results matrix with GSE ID, Title, Organism, Sample count, and Submission date."])
    add_card(slide7, Inches(6.8), Inches(4.6), Inches(5.7), Inches(2.1), "🔧 Raw Data & Trace", ["• Raw JSON response inspection from NCBI E-Utilities API calls.", "• Live reasoning trace panel showing agent scratchpad & tool outputs."])

    slide7.notes_slide.notes_text_frame.text = "Here is a demonstration of the visual dashboard outputs generated across the four tabs: Summary Metrics, Visual Charts, Datasets Table, and Raw API Inspection."


    # ── SLIDE 8: Hardware Optimization ────────────────────────────────────────
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide8, "Hardware & Performance", "Local GPU & Hardware Optimization", "High-efficiency local LLM inference configured for modern AI PCs.")

    add_card(slide8, Inches(0.8), Inches(2.2), Inches(3.6), Inches(4.5), 
             "💻 Workstation Specs", 
             ["• Target OS: Ubuntu 24.04 LTS",
              "• GPU: 8 GB VRAM GPU",
              "• RAM: 16+ GB RAM",
              "• CUDA: 12.1+"])

    add_card(slide8, Inches(4.8), Inches(2.2), Inches(3.6), Inches(4.5), 
             "🧠 Quantized LLM Profiles", 
             ["• medgemma:4b (~3.5 GB VRAM)",
              "  Default bio-focused model",
              "• gemma2:9b (~5.5 GB VRAM)",
              "  High-reasoning upgrade",
              "• llama3.1:8b (~4.9 GB VRAM)",
              "  General purpose alternative"])

    add_card(slide8, Inches(8.8), Inches(2.2), Inches(3.6), Inches(4.5), 
             "⚡ Batching & Limits", 
             ["• Embedding batch size: 512 (prevents CUDA OOM).",
              "• NCBI E-utils API: 10 req/s (with free API key).",
              "• Max dataset limit: Scalable up to 10,000 results per query."])

    slide8.notes_slide.notes_text_frame.text = "The toolkit is hardware-optimized for consumer and workstation GPUs with 8 GB VRAM. It uses medgemma:4b requiring only 3.5 GB VRAM with 512 embedding batch size."


    # ── SLIDE 9: Setup & Deployment ───────────────────────────────────────────
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide9, "Developer Guide", "Quick Start & Setup Workflow", "Simple 5-step terminal deployment for local development and research.")

    code_lines = [
        "# 1. One-time setup script (CUDA, Python venv, dependencies)",
        "chmod +x scripts/setup_ubuntu.sh",
        "./scripts/setup_ubuntu.sh",
        "",
        "# 2. Configure environment (.env)",
        "cp .env.example .env   # Set NCBI_EMAIL='your_email@example.com'",
        "",
        "# 3. Pull local GPU LLM model",
        "ollama pull medgemma:4b",
        "",
        "# 4. Build semantic ontology vector index",
        "python scripts/build_index.py",
        "",
        "# 5. Launch the interactive dashboard",
        "./venv/bin/python -m streamlit run app.py"
    ]

    add_card(slide9, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5), "Terminal Deployment Commands", code_lines, COLOR_CYAN)

    slide9.notes_slide.notes_text_frame.text = "Setting up the toolkit on Ubuntu is quick and automated. Run setup_ubuntu.sh, configure your .env file with your NCBI email, pull medgemma:4b via Ollama, build the vector index, and launch Streamlit."


    # ── SLIDE 10: Future Roadmap & Summary ────────────────────────────────────
    slide10 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide10, "Roadmap & Conclusion", "Future Vision & Summary", "Empowering bioinformaticians with self-healing autonomous research intelligence.")

    add_card(slide10, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.0), 
             "🚀 Key Achievements", 
             ["✓ 100% resilient 3-phase hybrid execution flow.",
              "✓ Sub-second statistics aggregation for up to 10,000 datasets.",
              "✓ Offline local GPU AI inference with medgemma:4b.",
              "✓ FAISS RAG semantic term normalization."], 
             COLOR_TEAL)

    add_card(slide10, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.0), 
             "🔮 Future Development Roadmap", 
             ["→ Spatial Transcriptomics & Single-Cell h5ad file parsing.",
              "→ Multi-Agent peer review consensus reports.",
              "→ Automated differential gene expression cross-comparison.",
              "→ PDF / LaTeX automated report generation."], 
             COLOR_CYAN)

    # Thank you footer box
    tb_ty = slide10.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8))
    p_ty = tb_ty.text_frame.paragraphs[0]
    p_ty.text = "Thank You! — Questions & Discussion"
    p_ty.alignment = PP_ALIGN.CENTER
    p_ty.font.size = Pt(20)
    p_ty.font.bold = True
    p_ty.font.color.rgb = COLOR_TEAL

    slide10.notes_slide.notes_text_frame.text = "In conclusion, the GEO Meta-Analysis Toolkit unlocks automated, self-healing biomedical data exploration. Future roadmap items include single-cell spatial transcriptomics integrations, multi-agent debates, and direct CSV/Excel export. Thank you! Any questions?"

    output_path = "presentation.pptx"
    prs.save(output_path)
    print(f"Successfully generated {output_path}")

if __name__ == "__main__":
    create_presentation()
